from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(255, 255, 255)
CARD = RGBColor(247, 249, 252)
ACCENT = RGBColor(43, 108, 176)
TEXT = RGBColor(26, 32, 44)
MUTED = RGBColor(74, 85, 104)
IMG_BG = RGBColor(237, 242, 247)
BORDER = RGBColor(203, 213, 224)

slides = [
    {
        "title": "MAFIA AI",
        "subtitle": "Учебный проект: изучение AI-технологий на практическом кейсе",
        "points": [
            "Главный результат проекта — не только продукт, а понимание технологий.",
            "Игра «Мафия» используется как удобный сценарий для проверки гипотез.",
        ],
        "img_label": "Титульная иллюстрация: команда + исследовательский процесс",
        "prompt": "clean white background, educational poster style, students presenting a software project, laptop with simple dashboard mockup, collaborative mood, minimal flat illustration, soft blue-gray accents, no neon, no futuristic effects, no text, 16:9",
        "speech": "В этой презентации мы делаем акцент на том, какие технологии изучили и как применили их на практике.",
    },
    {
        "title": "Зачем этот проект",
        "subtitle": "Проект как учебная площадка по AI и разработке",
        "points": [
            "Нам нужен был реальный кейс, где есть видео, аудио и логика принятия решений.",
            "Поэтому выбрали «Мафию»: тут можно изучать сразу несколько технологий.",
        ],
        "img_label": "Проблемная ситуация и учебная мотивация",
        "prompt": "simple infographic on white background, board game moderation challenges illustrated with icons: timer, speaking order, fairness, and a student notebook with AI labels, minimal style, academic clean layout, soft blue accents, no neon, no text, 16:9",
        "speech": "Мы выбрали проект так, чтобы каждая часть давала новый технический навык, а не просто красивый интерфейс.",
    },
    {
        "title": "Учебные цели",
        "subtitle": "Что именно мы хотели освоить",
        "points": [
            "Компьютерное зрение: лица, положение объектов, обработка потока с камеры.",
            "Обработка звука: идентификация голоса и распознавание речи.",
            "Веб-разработка: API, realtime-обмен, frontend/backend архитектура.",
        ],
        "img_label": "Дерево компетенций проекта",
        "prompt": "clean white competency map diagram, three branches: computer vision, audio processing, web architecture, each branch with simple icons and connectors, academic presentation style, minimal vector graphics, no neon, no futuristic effects, no text, 16:9",
        "speech": "На старте мы зафиксировали не только функциональные, но и учебные цели — это помогало оценивать прогресс.",
    },
    {
        "title": "Технологический стек",
        "subtitle": "Какие инструменты мы применили",
        "points": [
            "Frontend: React + TypeScript.",
            "Backend: FastAPI + WebSocket + REST API.",
            "AI-модули: распознавание лиц, голоса, речи, детекция стола.",
        ],
        "img_label": "Карта технологий и связей",
        "prompt": "clean architecture map on white background with labeled blocks style but without text rendered, web frontend block, backend api block, ai modules block, data storage block, arrows between blocks, corporate training infographic style, no neon, 16:9",
        "speech": "Важно, что мы работали не с одной библиотекой, а со связкой технологий, как в реальном проекте.",
    },
    {
        "title": "Как мы изучали Computer Vision",
        "subtitle": "От изображения до полезного сигнала",
        "points": [
            "Разобрались с видеопотоком, частотой кадров и предобработкой изображения.",
            "Освоили подход “детекция → признаки → сравнение”.",
        ],
        "img_label": "Пайплайн обработки видеокадра",
        "prompt": "educational computer vision pipeline diagram on white background: camera frame input, face detection box, feature extraction vectors, matching stage, output identity tag icon, clean minimal infographic, blue-gray palette, no neon, 16:9",
        "speech": "Основная мысль: система не “угадывает по фото”, а сравнивает числовые признаки, которые извлекаются из кадра.",
    },
    {
        "title": "Распознавание лиц (CompreFace)",
        "subtitle": "Что мы поняли на практике",
        "points": [
            "Качество зависит от освещения, ракурса и стабильности кадра.",
            "Нужны пороги уверенности и проверка ошибок, чтобы снизить ложные совпадения.",
        ],
        "img_label": "Схема face enrollment и face identify",
        "prompt": "white background process diagram, two branches: face enrollment and face identification, quality checks, confidence threshold concept, clean business infographic look, simple icons, no neon or futuristic effects, no text, 16:9",
        "speech": "Мы увидели, что алгоритм — это только часть задачи. Не меньше важны параметры, пороги и качество входных данных.",
    },
    {
        "title": "Работа со звуком и голосом",
        "subtitle": "Два слоя: кто говорит и что говорит",
        "points": [
            "Speaker ID: сопоставление голосового профиля игрока.",
            "ASR: преобразование речи в текст для логов и анализа.",
        ],
        "img_label": "Аудиопайплайн: микрофон → анализ → результат",
        "prompt": "clean white audio pipeline infographic, microphone input, waveform chunks, speaker profile matching branch, speech-to-text branch, merged output card, minimal academic style, soft blue and gray colors, no neon, 16:9",
        "speech": "Мы разделили задачу на две: идентификация спикера и распознавание текста, потому что это разные алгоритмы и разные ошибки.",
    },
    {
        "title": "Realtime и задержки",
        "subtitle": "Как мы балансировали скорость и качество",
        "points": [
            "Проблема: качественная транскрипция может увеличивать задержку.",
            "Решение: оптимизация чанков, очередей и таймаутов.",
        ],
        "img_label": "Баланс latency vs quality",
        "prompt": "white background engineering tradeoff chart illustration, balance scale between latency and accuracy, streaming queue and chunk timing icons, clean technical infographic style, no neon, no futuristic, no text, 16:9",
        "speech": "Один из главных технических уроков — всегда есть компромисс между скоростью и точностью.",
    },
    {
        "title": "Определение стола",
        "subtitle": "Зачем нужна геометрия сцены",
        "points": [
            "Система должна понимать рабочую область игры.",
            "Используем автоопределение и ручную корректировку контура.",
        ],
        "img_label": "Камера с контуром стола",
        "prompt": "clean white ui mockup style image, camera frame with table polygon overlay, auto detect button concept and manual point editing concept, minimal design, educational software screenshot style, no neon, 16:9",
        "speech": "Этот модуль показал, что даже простая геометрия сильно влияет на корректность всей системы.",
    },
    {
        "title": "Архитектура проекта",
        "subtitle": "Почему разбили систему на слои",
        "points": [
            "Frontend, Backend и AI-модули разделены, чтобы проще тестировать и развивать.",
            "Такой подход делает проект более поддерживаемым.",
        ],
        "img_label": "Слоистая архитектурная схема",
        "prompt": "clean layered architecture infographic on white background, presentation layer, api layer, application logic layer, infrastructure layer, arrows down, minimal flat style, soft neutral palette, no neon, no futuristic effects, no text, 16:9",
        "speech": "Архитектурное разделение помогло нам не запутаться и быстрее находить ошибки в конкретных частях системы.",
    },
    {
        "title": "Эксперименты и отладка",
        "subtitle": "Что пришлось решать по ходу разработки",
        "points": [
            "Конфликты портов, таймауты, шум в аудио, ложные распознавания.",
            "Мы научились диагностике: логи, тестовые роуты, итеративные улучшения.",
        ],
        "img_label": "Цикл: гипотеза → тест → лог → исправление",
        "prompt": "white background engineering iteration cycle diagram: hypothesis, experiment, logs, fix, retest, simple arrows in loop, notebook and terminal icons, clean educational infographic, no neon, 16:9",
        "speech": "Здесь мы получили самый важный практический опыт: как системно разбирать и исправлять технические проблемы.",
    },
    {
        "title": "Что мы получили как учебный результат",
        "subtitle": "Ключевые навыки команды",
        "points": [
            "Интеграция AI-модулей в реальное приложение.",
            "Работа с realtime-системами и пользовательскими сценариями.",
            "Командная разработка и инженерная дисциплина.",
        ],
        "img_label": "Итоговые навыки и компетенции",
        "prompt": "clean white summary infographic with skill icons: computer vision, speech processing, backend api, frontend ui, debugging, teamwork, educational style, minimal layout, no neon, no futuristic effects, no text, 16:9",
        "speech": "Если кратко: этот проект дал нам не только демо-продукт, а полноценный набор практических инженерных навыков.",
    },
    {
        "title": "Вывод",
        "subtitle": "Проект как способ понять технологии через практику",
        "points": [
            "Главная ценность: мы глубоко разобрались в современных технологиях AI и веб-разработки.",
            "Продукт — это демонстрация того, чему мы научились.",
        ],
        "img_label": "Финальный нейтральный кадр для завершения",
        "prompt": "clean white closing slide illustration, students presenting project outcomes to audience, simple charts and mockups on screen, positive academic atmosphere, minimalist style, no neon, no futuristic effects, no text, 16:9",
        "speech": "Для нас это в первую очередь учебный исследовательский проект, который помог перевести теорию в работающую систему.",
    },
]


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, title, subtitle):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(9.8), Inches(1.0))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(34)
    run.font.bold = True
    run.font.color.rgb = TEXT

    sub_box = slide.shapes.add_textbox(Inches(0.62), Inches(1.2), Inches(12.0), Inches(0.65))
    tf2 = sub_box.text_frame
    tf2.clear()
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = subtitle
    run2.font.size = Pt(16)
    run2.font.color.rgb = MUTED


def add_points(slide, points):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.9), Inches(5.3), Inches(2.4))
    box.fill.solid()
    box.fill.fore_color.rgb = CARD
    box.line.color.rgb = BORDER
    box.line.width = Pt(1.1)

    tf = box.text_frame
    tf.clear()
    for i, txt in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT
        p.level = 0
        p.space_after = Pt(8)


def add_image_placeholder(slide, label):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.05), Inches(1.9), Inches(6.55), Inches(4.95))
    shp.fill.solid()
    shp.fill.fore_color.rgb = IMG_BG
    shp.line.color.rgb = BORDER
    shp.line.width = Pt(1.3)

    tf = shp.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "МЕСТО ДЛЯ ИЗОБРАЖЕНИЯ"
    run.font.size = Pt(17)
    run.font.bold = True
    run.font.color.rgb = ACCENT

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = label
    p2.font.size = Pt(12)
    p2.font.color.rgb = MUTED


for idx, data in enumerate(slides, start=1):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_title(slide, f"{idx}. {data['title']}", data["subtitle"])
    add_points(slide, data["points"])
    add_image_placeholder(slide, data["img_label"])

    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = (
        "ПРОМПТ ДЛЯ КАРТИНКИ:\n"
        + data["prompt"]
        + "\n\nЧТО ГОВОРИТЬ НА СЛАЙДЕ (20-30 сек):\n"
        + data["speech"]
        + "\n"
    )

out = "presentation/Mafia_AI_tech_focus_clean_white_google_slides_ready.pptx"
prs.save(out)
print(out)
