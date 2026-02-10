from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Black presentation theme (calm, no neon)
BG = RGBColor(14, 17, 22)
CARD = RGBColor(25, 30, 38)
ACCENT = RGBColor(126, 162, 209)
TEXT = RGBColor(237, 242, 247)
MUTED = RGBColor(170, 180, 196)
IMG_BG = RGBColor(31, 38, 49)
BORDER = RGBColor(64, 74, 90)

slides = [
    {
        "title": "MAFIA AI",
        "subtitle": "Учебный проект: как мы разобрались в AI-технологиях на реальном кейсе",
        "points": [
            "Главный результат — освоение технологий, а не только готовый продукт.",
            "Игра «Мафия» стала удобной лабораторией для видео, аудио и backend/frontend интеграции.",
        ],
        "img_label": "Команда и общий образ проекта",
        "prompt": "dark neutral background, students presenting software project, laptop with simple dashboard mockup, realistic educational atmosphere, minimal clean composition, no neon, no sci-fi, no glowing holograms, no text, 16:9",
        "speech": "На этом слайде важно задать правильный фокус всей презентации. Скажите, что наш проект называется Mafia AI, но это в первую очередь учебный инженерный проект. Мы не ставили цель просто сделать красивое приложение, мы ставили цель понять, как работают современные технологии: компьютерное зрение, обработка звука, распознавание речи, realtime-системы и архитектура веб-приложений. Игра «Мафия» стала для нас удобным сценарием, где все эти технологии можно проверить в одной связке. Дальше в презентации мы будем показывать именно технологический путь: какие задачи мы решали, какие решения пробовали, с какими проблемами сталкивались, и что это дало нам как разработчикам.",
    },
    {
        "title": "Почему мы выбрали именно такой кейс",
        "subtitle": "Проект должен был быть одновременно понятным и технически насыщенным",
        "points": [
            "Нужен сценарий, где есть видео, голос, события в реальном времени и логика правил.",
            "«Мафия» подходит идеально: участники, очередь речи, контроль нарушений, этапы игры.",
        ],
        "img_label": "Проблематика и учебная мотивация",
        "prompt": "realistic board game discussion scene in dark room with neutral lighting, people around table, timer and notes visible, educational documentary style, clean composition, no neon, no futuristic elements, no text, 16:9",
        "speech": "Здесь поясните выбор темы. Мы специально искали задачу, которая была бы понятна любому человеку, но при этом технически многослойна. В обычной игре «Мафия» есть много моментов, которые сложно контролировать вручную: кто говорил, когда говорил, кто нарушил порядок, как фиксировать этапы. Для нас это означало, что в проекте одновременно появятся и видеоаналитика, и аудиоаналитика, и backend-логика, и интерфейс для пользователя. То есть одна игра дала нам сразу несколько учебных направлений. Это важно подчеркнуть: мы выбрали тему не потому, что хотели только развлечение, а потому что она хорошо моделирует реальные инженерные задачи, где нужно объединять разные технологии в единую систему.",
    },
    {
        "title": "Наши учебные цели",
        "subtitle": "Какие технологические компетенции мы целенаправленно прокачивали",
        "points": [
            "Computer Vision: детекция и идентификация лиц на видеопотоке.",
            "Audio/NLP: идентификация говорящего и транскрипция речи.",
            "Web Architecture: API, WebSocket, frontend/backend взаимодействие.",
        ],
        "img_label": "Карта компетенций проекта",
        "prompt": "clean academic infographic on dark gray background, three branches: computer vision, audio processing, web architecture, simple icons and connectors, minimal style, no neon, no futuristic style, no text, 16:9",
        "speech": "На этом слайде перечислите не функции приложения, а именно навыки и технологии, которые мы хотели изучить. Первая группа — компьютерное зрение: как получать кадры с камеры, как находить на них лица, как сравнивать лицо с базой. Вторая — работа с аудио и речью: как отделять полезную речь от шума, как определять говорящего, как получать текст из аудио. Третья — архитектурная часть: как построить frontend и backend, как сделать API, как организовать realtime-обновления через WebSocket. Можно сказать, что у нас был не просто список «что должно работать», а учебный план «что мы должны понять». И итог проекта мы оцениваем именно по тому, насколько глубоко разобрались в этих технологических блоках.",
    },
    {
        "title": "Технологический стек",
        "subtitle": "Из каких инструментов собрана система",
        "points": [
            "Frontend: React + TypeScript (интерфейс и сценарии настройки).",
            "Backend: FastAPI (REST + WebSocket, бизнес-логика, маршруты).",
            "AI-компоненты: распознавание лиц, голосов, речи и определение стола.",
        ],
        "img_label": "Блок-схема стека",
        "prompt": "software architecture blocks on dark matte background, frontend backend ai modules storage connected with arrows, minimal professional infographic style, no neon, no futuristic, no text, 16:9",
        "speech": "Здесь проговорите стек простыми словами. Интерфейс сделан на React и TypeScript — это то, с чем работает пользователь. Серверная часть на FastAPI — она принимает запросы, хранит данные, управляет игровыми этапами и раздает realtime-события. Отдельно подключены AI-компоненты: распознавание лиц, распознавание и идентификация речи, работа с геометрией стола. Важно отметить, что основной учебный эффект дал именно процесс интеграции: каждая технология по отдельности может быть понятна, но когда их нужно объединить в рабочий pipeline, появляются настоящие инженерные задачи — форматы данных, задержки, устойчивость, обработка ошибок.",
    },
    {
        "title": "Как мы разбирались с Computer Vision",
        "subtitle": "От кадра камеры до идентифицированного игрока",
        "points": [
            "Кадр с камеры → поиск лица → вектор признаков → сравнение с базой.",
            "Качество зависит от света, ракурса и стабильности изображения.",
        ],
        "img_label": "Пайплайн обработки видеокадра",
        "prompt": "computer vision process illustration on dark neutral background, camera frame with detected face box, feature vector abstraction, matching against profile gallery, clean educational style, realistic and minimal, no neon, no sci-fi, no text, 16:9",
        "speech": "Здесь объясняйте без математики. Система берет кадр с камеры, находит лицо, затем превращает его в числовое представление — условно «отпечаток признаков». Дальше этот отпечаток сравнивается с уже зарегистрированными профилями игроков. Если совпадение достаточно уверенное, система считает, что это конкретный игрок. Отдельно подчеркните практический урок: точность сильно зависит от условий съемки. Мы увидели, что для реальной работы важны не только алгоритмы, но и качественная регистрация, освещение и адекватные пороги уверенности. То есть Computer Vision — это не магия одной кнопки, это комбинация модели, параметров и условий среды.",
    },
    {
        "title": "Распознавание лиц через CompreFace",
        "subtitle": "Интеграция внешнего движка и инженерные выводы",
        "points": [
            "Использовали отдельный сервис распознавания лиц через API.",
            "Научились настраивать пороги и корректно обрабатывать ошибки распознавания.",
        ],
        "img_label": "Схема: приложение ↔ сервис распознавания лиц",
        "prompt": "backend service integration diagram on dark gray background, app server communicating with face recognition service, request/response arrows, quality threshold concept, clean enterprise style, no neon, no futuristic effects, no text, 16:9",
        "speech": "Этот слайд важен как пример реальной интеграции. Мы не просто использовали локальный скрипт, а работали с отдельным сервисом распознавания лиц через API. Это дало опыт, близкий к промышленной разработке: как поднимать сервис, как передавать данные, как учитывать таймауты, как обрабатывать случаи, когда лицо не найдено. Также мы на практике поработали с порогами уверенности: слишком высокий порог дает пропуски, слишком низкий — ложные совпадения. Главный учебный вывод: инженер должен уметь не только «включить технологию», но и подобрать рабочие параметры под реальные условия и контролировать качество результата.",
    },
    {
        "title": "Голос и речь: два разных уровня",
        "subtitle": "Кто говорит и что именно сказано",
        "points": [
            "Speaker ID: определяем личность говорящего по голосовому профилю.",
            "ASR: распознаем текст речи для логов и анализа игрового процесса.",
        ],
        "img_label": "Аудиопайплайн и разветвление задач",
        "prompt": "audio processing pipeline on dark neutral background, microphone input split into speaker identification branch and speech-to-text branch, merged output log, clean educational diagram style, no neon, no sci-fi, no text, 16:9",
        "speech": "На этом слайде подчеркните, что «распознавание голоса» и «распознавание речи» — это разные задачи. Первая отвечает на вопрос «кто говорит», вторая — «что сказано». Мы обучали голосовые профили игроков и сравнивали текущий аудиофрагмент с этими профилями. Параллельно запускали транскрипцию, чтобы получить текстовую ленту событий. Практический урок здесь в том, что аудио очень чувствительно к шуму и качеству микрофона. Поэтому пришлось работать с предобработкой сигнала, таймаутами и размером аудиочанков. Это дало нам понимание, как строить устойчивый аудиопайплайн для realtime-сценариев.",
    },
    {
        "title": "Realtime и задержка",
        "subtitle": "Компромисс по двум каналам: голос и лица",
        "points": [
            "Аудио: меньшие чанки дают быстрее лог, но снижают стабильность транскрипции.",
            "Лица: более точная модель и низкие пороги повышают нагрузку и latency.",
            "Подбирали общий баланс под железо: FPS, таймауты, размер очередей и confidence-пороги.",
        ],
        "img_label": "Trade-off: latency vs quality для голоса и лиц",
        "prompt": "clean engineering infographic on dark graphite background, centered two-column tradeoff layout, left column for audio pipeline latency and right column for face recognition latency, balanced scales and queue icons, stopwatch icon, cpu load icon, confidence badge icon, chunked stream blocks flowing into processing nodes, subtle arrows showing compromise between speed and quality for both channels, minimal academic style, restrained glassmorphism cards, no neon, no futuristic style, no text, 16:9",
        "speech": "На этом слайде важно сказать, что компромисс в realtime у нас был не только по голосу, но и по лицам. По аудио: если брать маленькие чанки, лог появляется быстрее, но растет риск ошибок транскрипции. По лицам: если использовать более тяжелые модели или слишком строгие пороги, качество может расти, но задержка и нагрузка увеличиваются. Поэтому мы настраивали систему как единый контур: частоту кадров, размер очередей, таймауты и confidence-пороги. Главный вывод — в realtime-проекте нельзя оптимизировать только один модуль, нужно балансировать весь пайплайн под конкретное железо и требования к удобству интерфейса.",
    },
    {
        "title": "Транскрипция жестов",
        "subtitle": "Как визуальные сигналы превращаются в текстовые логи",
        "points": [
            "Детектим жесты рук в видеопотоке и нормализуем их в токены.",
            "Собираем устойчивую последовательность токенов и переводим в фразу.",
            "Добавляем результат в общий лог как отдельный тип события.",
        ],
        "img_label": "Пайплайн: жесты → токены → текст",
        "prompt": "clean academic process diagram on dark neutral background, camera frame with two hands and landmark points, gesture tokens pipeline (numbers and symbolic gestures) flowing into normalization module then phrase builder module then speech log output card, clear sequence from visual gesture to readable phrase, minimal enterprise infographic, restrained glass cards, no neon, no sci-fi, no futuristic style, no text, 16:9",
        "speech": "Этот слайд объясняет новый модуль транскрипции жестов. Сначала система в потоке выделяет руки и определяет жесты: цифры и условные символы. Дальше эти жесты нормализуются в токены, после чего собираются в устойчивую последовательность, чтобы убрать случайные колебания распознавания. Затем последовательность переводится в читаемую фразу и записывается в лог как отдельный тип события. Для нас это важный учебный шаг: мы сделали связку Computer Vision и семантической интерпретации, то есть перешли от детекции отдельных жестов к осмысленному текстовому результату.",
    },
    {
        "title": "Определение стола",
        "subtitle": "Геометрия сцены как часть AI-контекста",
        "points": [
            "Система должна понимать, где находится игровая зона.",
            "Добавили автоопределение и ручную разметку контура для надежности.",
        ],
        "img_label": "Кадр с контуром стола",
        "prompt": "camera frame with table polygon overlay on dark neutral interface, auto-detect and manual point editing concepts, realistic room, clean UI mockup style, no neon, no futuristic effects, no text, 16:9",
        "speech": "Здесь поясните идею контекста сцены. Если система не знает, где стол, ей сложнее корректно интерпретировать действия игроков. Поэтому мы сделали модуль калибровки: либо автоматическое определение контура, либо ручная корректировка, если автоматике нужна помощь. Этот шаг уменьшает количество ошибок в дальнейшем анализе. Учебный вывод: даже относительно простой геометрический модуль может заметно повысить качество всей системы. Это хороший пример того, как не самые «громкие» алгоритмы дают большой практический эффект в продукте.",
    },
    {
        "title": "Архитектура проекта",
        "subtitle": "Почему нам было важно разделить систему на слои",
        "points": [
            "Frontend — взаимодействие с пользователем.",
            "Backend — логика, API и orchestration модулей.",
            "AI-блоки — специализированная обработка данных.",
        ],
        "img_label": "Слоистая архитектурная диаграмма",
        "prompt": "layered software architecture diagram on dark gray background, presentation layer, api/backend layer, application logic, infrastructure and ai modules, clean simple arrows, enterprise educational style, no neon, no futuristic, no text, 16:9",
        "speech": "На этом слайде покажите, что мы думали не только о фичах, но и о структуре кода. Мы разделили проект на слои, чтобы каждый блок отвечал за свою задачу: интерфейс, серверная логика, AI-модули. Это облегчает отладку и развитие: если проблема в распознавании речи, не нужно трогать UI; если меняется интерфейс, не ломается логика распознавания. Для школьного проекта это особенно сильный момент, потому что показывает инженерное мышление. Подчеркните, что такая архитектура помогает работать командой и постепенно улучшать систему, не переписывая всё с нуля.",
    },
    {
        "title": "Эксперименты и отладка",
        "subtitle": "Как мы решали реальные технические проблемы",
        "points": [
            "Столкнулись с портами, таймаутами, шумом, неточными срабатываниями.",
            "Использовали логи, тестовые маршруты и итеративные исправления.",
        ],
        "img_label": "Цикл инженерной итерации",
        "prompt": "software debugging workflow on dark neutral background, cycle of hypothesis test logs fix retest, laptop terminal and notebook icons, clean practical style, no neon, no futuristic, no text, 16:9",
        "speech": "Этот слайд посвятите процессу работы. Расскажите, что проект развивался через цикл: гипотеза, тест, анализ логов, исправление, повторный тест. Мы регулярно сталкивались с конкретными проблемами: занятые порты, таймауты при тяжелой обработке, задержка в логах, ложные распознавания. И именно здесь мы получили самый ценный опыт — как системно диагностировать проблему и подтверждать, что исправление действительно работает. Можно сказать, что именно этап отладки превратил проект из «прототипа» в инженерный учебный кейс. Этот опыт напрямую переносится на любые будущие реальные разработки.",
    },
    {
        "title": "Главный учебный результат",
        "subtitle": "Какие навыки мы реально получили",
        "points": [
            "Понимание полного цикла AI-фичи: от данных до интерфейса.",
            "Опыт интеграции сервисов, realtime-коммуникации и отладки.",
            "Умение принимать компромиссные технические решения.",
        ],
        "img_label": "Итоговая карта навыков команды",
        "prompt": "skills summary infographic on dark matte background, icons for computer vision, speech processing, backend api, frontend ui, realtime systems, debugging, teamwork, clean academic style, no neon, no futuristic effects, no text, 16:9",
        "speech": "На этом слайде подведите итог не в формате «что умеет продукт», а в формате «что теперь умеем мы». Мы прошли полный путь: формулировка задачи, выбор инструментов, интеграция, диагностика ошибок, улучшение производительности, пользовательский сценарий. Мы получили не фрагментарные знания, а системное понимание того, как строятся современные AI-функции в приложениях. Это и есть главный образовательный результат проекта. И уже как следствие появился рабочий продукт, который можно показать и протестировать. То есть продукт — это доказательство освоенных технологий, а не самоцель.",
    },
    {
        "title": "Вывод",
        "subtitle": "Проект как практическое изучение технологий",
        "points": [
            "Mafia AI — это учебная лаборатория, собранная в формате работающего приложения.",
            "Главная ценность — технологическая глубина, которую мы получили в процессе разработки.",
        ],
        "img_label": "Финальный кадр для завершения",
        "prompt": "students presenting project conclusions in classroom, dark neutral background, calm realistic lighting, screen with simple charts and app mockups, professional educational mood, no neon, no futuristic effects, no text, 16:9",
        "speech": "В финале еще раз закрепите главный тезис всей презентации. Наш проект важен в первую очередь как путь изучения технологий: компьютерного зрения, аудиообработки, веб-архитектуры и realtime-систем. Мы получили опыт не только в том, как «сделать, чтобы работало», но и в том, как анализировать качество, устранять ошибки и принимать инженерные компромиссы. Поэтому этот проект для нас — не просто демонстрация приложения, а доказательство того, что мы умеем работать с современным технологическим стеком на практике. После этого можно предложить показать короткое демо или перейти к вопросам.",
    },
]


def set_slide_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, title, subtitle):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(9.9), Inches(1.0))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(34)
    run.font.bold = True
    run.font.color.rgb = TEXT

    sub_box = slide.shapes.add_textbox(Inches(0.62), Inches(1.2), Inches(12.0), Inches(0.65))
    tf2 = sub_box.text_frame
    tf2.clear()
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = subtitle
    run2.font.size = Pt(16)
    run2.font.color.rgb = MUTED


def add_points(slide, points):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.9), Inches(5.35), Inches(2.55))
    box.fill.solid()
    box.fill.fore_color.rgb = CARD
    box.line.color.rgb = BORDER
    box.line.width = Pt(1.1)

    tf = box.text_frame
    tf.clear()
    for i, txt in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT
        p.level = 0
        p.space_after = Pt(7)


def add_image_placeholder(slide, label):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.05), Inches(1.9), Inches(6.55), Inches(4.95))
    shp.fill.solid()
    shp.fill.fore_color.rgb = IMG_BG
    shp.line.color.rgb = BORDER
    shp.line.width = Pt(1.3)

    tf = shp.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "МЕСТО ДЛЯ ИЗОБРАЖЕНИЯ"
    run.font.size = Pt(17)
    run.font.bold = True
    run.font.color.rgb = ACCENT

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = label
    p2.font.size = Pt(12)
    p2.font.color.rgb = MUTED


for i, data in enumerate(slides, start=1):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, f"{i}. {data['title']}", data['subtitle'])
    add_points(slide, data['points'])
    add_image_placeholder(slide, data['img_label'])

    notes = slide.notes_slide.notes_text_frame
    notes.clear()
    notes.text = (
        "ПРОМПТ ДЛЯ КАРТИНКИ:\n"
        + data['prompt']
        + "\n\nПОДРОБНЫЙ ТЕКСТ ВЫСТУПЛЕНИЯ НА СЛАЙДЕ:\n"
        + data['speech']
        + "\n"
    )

out = "presentation/Mafia_AI_tech_focus_black_detailed_notes_google_slides_ready.pptx"
prs.save(out)
print(out)
